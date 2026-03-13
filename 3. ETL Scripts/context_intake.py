"""
FintastIQ Client Context Intake Module
Captures business context before data processing — feeds smarter LLM decisions
throughout the pipeline.

3-Tier Hybrid Approach:
  Tier 1: Required structured fields (company, industry, model, revenue, product, goal)
  Tier 2: LLM-generated follow-up questions based on Tier 1 answers
  Tier 3: Open-ended context + deck/doc extraction

Input paths:
  A) FintastIQ team fills from proposal deck, kickoff notes, or manual entry
  B) Client self-serves via portal
  C) LLM extracts from uploaded .pptx/.docx

Output: client_context.json — injected into every pipeline stage
"""

import json, os, re, time
from datetime import datetime

import httpx, anthropic

# ── CONFIG ──
API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
MODEL = "claude-sonnet-4-20250514"

def get_client():
    return anthropic.Anthropic(api_key=API_KEY, http_client=httpx.Client(verify=False, timeout=120.0))

def llm_call(prompt, max_tokens=4096):
    client = get_client()
    resp = client.messages.create(model=MODEL, max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}])
    return resp.content[0].text, resp.usage

def extract_json(text, prefer_array=False):
    first_char = '[' if prefer_array else '{'
    second_char = '{' if prefer_array else '['
    for ch in [first_char, second_char]:
        idx = text.find(ch)
        if idx == -1:
            continue
        depth = 0
        close = ']' if ch == '[' else '}'
        for i in range(idx, len(text)):
            if text[i] == ch: depth += 1
            elif text[i] == close: depth -= 1
            if depth == 0:
                try:
                    return json.loads(text[idx:i+1])
                except:
                    break
    return None


# ══════════════════════════════════════════════════════════════
# TIER 1: REQUIRED STRUCTURED FIELDS
# ══════════════════════════════════════════════════════════════

TIER1_SCHEMA = {
    "company_name": {
        "label": "Company Name",
        "type": "text",
        "required": True,
        "hint": "Full legal or commonly used name",
        "example": "NPI Financial"
    },
    "industry": {
        "label": "Industry Vertical",
        "type": "text",
        "required": True,
        "hint": "Primary industry (e.g., B2B SaaS, Manufacturing, Healthcare, Financial Services)",
        "example": "B2B SaaS / Consulting"
    },
    "business_model": {
        "label": "Business Model",
        "type": "select",
        "required": True,
        "options": ["B2B", "B2C", "B2B2C", "Marketplace", "Hybrid"],
        "hint": "How does the company primarily sell?",
        "example": "B2B"
    },
    "revenue_range": {
        "label": "Approximate Annual Revenue",
        "type": "select",
        "required": True,
        "options": ["<$1M", "$1M-$10M", "$10M-$50M", "$50M-$250M", "$250M-$1B", "$1B+"],
        "hint": "Rough annual revenue range",
        "example": "$10M-$50M"
    },
    "product_type": {
        "label": "What They Sell",
        "type": "select_multi",
        "required": True,
        "options": ["Physical Products", "Software/SaaS", "Professional Services",
                    "Subscriptions", "Licensing", "Marketplace/Platform", "Mixed"],
        "hint": "Primary revenue drivers",
        "example": ["Professional Services", "Subscriptions"]
    },
    "engagement_goal": {
        "label": "Primary Engagement Goal",
        "type": "text",
        "required": True,
        "hint": "One-liner: what does the client want from this pricing diagnostic?",
        "example": "Identify pricing leakage and optimize win rates"
    }
}


def validate_tier1(context):
    """Validate that all required Tier 1 fields are present and non-empty."""
    missing = []
    for field, schema in TIER1_SCHEMA.items():
        if schema["required"]:
            val = context.get(field)
            if val is None or (isinstance(val, str) and val.strip() == ''):
                missing.append(field)
            elif isinstance(val, list) and len(val) == 0:
                missing.append(field)
    return missing


# ══════════════════════════════════════════════════════════════
# TIER 2: LLM-GENERATED FOLLOW-UP QUESTIONS
# ══════════════════════════════════════════════════════════════

def generate_tier2_questions(tier1_context):
    """Generate 3-5 targeted follow-up questions based on Tier 1 answers."""
    prompt = f"""You are an expert pricing consultant preparing for a client engagement.
Based on this client profile, generate 3-5 targeted follow-up questions that will help
us better understand their pricing structure and data when we process their files.

Client Profile:
- Company: {tier1_context.get('company_name', 'Unknown')}
- Industry: {tier1_context.get('industry', 'Unknown')}
- Business Model: {tier1_context.get('business_model', 'Unknown')}
- Revenue Range: {tier1_context.get('revenue_range', 'Unknown')}
- Product/Service: {json.dumps(tier1_context.get('product_type', []))}
- Engagement Goal: {tier1_context.get('engagement_goal', 'Unknown')}

Return a JSON object with:
{{
  "questions": [
    {{
      "id": "Q1",
      "question": "...",
      "why_it_matters": "Brief explanation of how this helps the data pipeline",
      "answer_type": "text" | "select" | "number",
      "options": ["..."] // only if answer_type is "select"
    }},
    ...
  ]
}}

Rules:
- Questions should be specific to THIS client's industry and model
- Focus on things that affect how we interpret their data
- Ask about pricing structure, contract types, discount authority, key metrics
- Don't ask things already covered in the profile above
- Keep questions concise and business-friendly (not technical)
- 3-5 questions maximum"""

    t0 = time.time()
    text, usage = llm_call(prompt, max_tokens=2048)
    result = extract_json(text)
    elapsed = time.time() - t0

    if result and 'questions' in result:
        print(f"  Generated {len(result['questions'])} follow-up questions ({elapsed:.1f}s)")
        return result['questions'], usage
    else:
        print(f"  Failed to generate questions — using defaults")
        return _default_questions(tier1_context), usage


def _default_questions(tier1_context):
    """Fallback questions if LLM generation fails."""
    bm = tier1_context.get('business_model', '')
    questions = [
        {"id": "Q1", "question": "How is your pricing structured? (e.g., per-seat, usage-based, project-based, tiered)",
         "why_it_matters": "Determines how we interpret revenue data and calculate unit economics",
         "answer_type": "text"},
        {"id": "Q2", "question": "What does a typical sales cycle look like? (length, stages, key decision points)",
         "why_it_matters": "Helps interpret deal pipeline and win/loss patterns",
         "answer_type": "text"},
        {"id": "Q3", "question": "Who has authority to approve discounts, and what are the typical discount ranges?",
         "why_it_matters": "Critical for discount governance and price leakage analyses",
         "answer_type": "text"},
    ]
    if bm in ('B2B', 'B2B2C'):
        questions.append(
            {"id": "Q4", "question": "What is your customer retention/renewal rate?",
             "why_it_matters": "Key input for cohort and churn analyses",
             "answer_type": "text"})
    return questions


# ══════════════════════════════════════════════════════════════
# TIER 3: OPEN CONTEXT + DECK/DOC EXTRACTION
# ══════════════════════════════════════════════════════════════

TIER3_PROMPTS = {
    "pricing_structure": {
        "label": "Pricing Structure",
        "prompt": "Describe your pricing structure in your own words — how do you charge customers?",
        "required": False
    },
    "known_issues": {
        "label": "Known Pricing Issues",
        "prompt": "What's broken or suboptimal about your current pricing? What do you suspect?",
        "required": False
    },
    "data_quirks": {
        "label": "Data Quirks",
        "prompt": "Anything unusual about your data we should know? (e.g., naming conventions, currency, fiscal year, data gaps)",
        "required": False
    },
    "key_terms": {
        "label": "Key Business Terms",
        "prompt": "Any company-specific terms or acronyms we'll see in the data? (e.g., 'TCV = Total Contract Value', 'MRR = Monthly Recurring Revenue')",
        "required": False
    }
}


def extract_context_from_text(raw_text, tier1_context=None):
    """Extract structured context from free-form text (meeting notes, emails, etc.)."""
    existing = json.dumps(tier1_context, indent=2) if tier1_context else "None yet"

    prompt = f"""You are extracting client context from unstructured text for a pricing analytics platform.

Existing context (if any):
{existing}

Text to extract from:
---
{raw_text[:8000]}
---

Extract any relevant business context and return a JSON object with these possible fields:
{{
  "company_name": "...",
  "industry": "...",
  "business_model": "B2B|B2C|B2B2C|Marketplace|Hybrid",
  "revenue_range": "<$1M|$1M-$10M|$10M-$50M|$50M-$250M|$250M-$1B|$1B+",
  "product_type": ["..."],
  "engagement_goal": "...",
  "pricing_structure": "...",
  "known_issues": "...",
  "data_quirks": "...",
  "key_terms": {{"TERM": "definition", ...}},
  "additional_context": "anything else relevant"
}}

Rules:
- Only include fields you can confidently extract from the text
- If a field is already in existing context AND the text doesn't contradict it, omit it
- Use exact values from the option lists where applicable
- For key_terms, extract any defined acronyms or business terminology
- Return only the JSON, no commentary"""

    text, usage = llm_call(prompt, max_tokens=2048)
    result = extract_json(text)
    return result, usage


def extract_context_from_deck(pptx_path):
    """Extract context from a PowerPoint presentation file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("  python-pptx not installed — run: pip install python-pptx")
        return None, None

    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_content.append(" | ".join(row_text))
        if slide_content:
            slides_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_content))

    if not slides_text:
        print("  No text content found in deck")
        return None, None

    full_text = "\n\n".join(slides_text)
    print(f"  Extracted text from {len(slides_text)} slides ({len(full_text):,} chars)")
    return extract_context_from_text(full_text)


def extract_context_from_docx(docx_path):
    """Extract context from a Word document."""
    try:
        from docx import Document
    except ImportError:
        print("  python-docx not installed — run: pip install python-docx")
        return None, None

    doc = Document(docx_path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                paragraphs.append(" | ".join(row_text))

    if not paragraphs:
        return None, None

    full_text = "\n".join(paragraphs)
    print(f"  Extracted text from docx ({len(paragraphs)} paragraphs, {len(full_text):,} chars)")
    return extract_context_from_text(full_text)


def extract_context_from_file(file_path, tier1_context=None):
    """Auto-detect file type and extract context."""
    ext = os.path.splitext(file_path)[1].lower()
    print(f"  Extracting context from {os.path.basename(file_path)} ({ext})...")

    if ext in ('.pptx', '.ppt'):
        return extract_context_from_deck(file_path)
    elif ext in ('.docx', '.doc'):
        return extract_context_from_docx(file_path)
    elif ext in ('.txt', '.md'):
        with open(file_path, 'r', errors='ignore') as f:
            text = f.read()
        return extract_context_from_text(text, tier1_context)
    elif ext == '.pdf':
        # Try pdfplumber or fallback
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join(page.extract_text() or '' for page in pdf.pages[:20])
            return extract_context_from_text(text, tier1_context)
        except ImportError:
            print("  pdfplumber not installed — cannot extract from PDF")
            return None, None
    else:
        print(f"  Unsupported file type: {ext}")
        return None, None


# ══════════════════════════════════════════════════════════════
# CROSS-VALIDATION: CONTEXT vs ACTUAL DATA
# ══════════════════════════════════════════════════════════════

def cross_validate_context(context, data_folder):
    """Cross-check context against actual uploaded data files."""
    import pandas as pd
    import glob

    print(f"\n  Cross-validating context against data in {os.path.basename(data_folder)}...")

    # Scan data files for column names and sample values
    files = glob.glob(os.path.join(data_folder, '**/*.xlsx'), recursive=True) + \
            glob.glob(os.path.join(data_folder, '**/*.xlsb'), recursive=True) + \
            glob.glob(os.path.join(data_folder, '**/*.csv'), recursive=True)

    all_columns = set()
    sample_values = {}
    file_count = 0
    total_rows = 0

    for f in files[:20]:  # Cap at 20 files to keep it fast
        ext = os.path.splitext(f)[1].lower()
        try:
            if ext == '.csv':
                df = pd.read_csv(f, nrows=100)
            elif ext == '.xlsb':
                df = pd.read_excel(f, engine='pyxlsb', nrows=100)
            else:
                xl = pd.ExcelFile(f)
                # Read first sheet only
                df = pd.read_excel(f, sheet_name=xl.sheet_names[0], nrows=100)
            file_count += 1
            total_rows += len(df)
            for col in df.columns:
                col_str = str(col).strip()
                if col_str and len(col_str) < 60:
                    all_columns.add(col_str)
                    # Grab a few unique non-null values
                    vals = df[col].dropna().unique()[:5]
                    sample_values[col_str] = [str(v)[:80] for v in vals]
        except Exception:
            continue

    if file_count == 0:
        return {"flags": [], "confidence": "low", "note": "No data files could be read"}

    # Build a summary of what we found in the data
    data_summary = {
        "files_scanned": file_count,
        "total_sample_rows": total_rows,
        "column_names": sorted(list(all_columns))[:100],
        "sample_columns": {k: v for k, v in list(sample_values.items())[:30]}
    }

    prompt = f"""You are validating client-provided context against their actual data files.

Client Context:
{json.dumps(context, indent=2, default=str)}

Data Summary (from scanning their uploaded files):
- Files scanned: {data_summary['files_scanned']}
- Column names found: {json.dumps(data_summary['column_names'][:60])}
- Sample values from key columns:
{json.dumps(data_summary['sample_columns'], indent=2, default=str)}

Check for:
1. CONTRADICTIONS: Does the data contradict the stated business model, industry, or product type?
   (e.g., client says "B2C retail" but data has contract_value, renewal_date → likely B2B)
2. MISSING CONTEXT: What important context is missing that the data suggests?
   (e.g., data has "currency" column → should know if multi-currency)
3. CONFIRMED: What context is confirmed by the data?
4. INFERRED: What additional context can you confidently infer from the data?

Return a JSON object:
{{
  "flags": [
    {{
      "type": "contradiction|missing|confirmed|inferred",
      "field": "field_name",
      "message": "What we found",
      "severity": "high|medium|low",
      "suggested_value": "..." // if we can suggest a correction
    }},
    ...
  ],
  "inferred_fields": {{
    // Any new fields we can confidently add to the context
    "field_name": "value",
    ...
  }},
  "confidence": "high|medium|low",
  "notes": "Brief overall assessment"
}}"""

    t0 = time.time()
    text, usage = llm_call(prompt, max_tokens=2048)
    result = extract_json(text)
    elapsed = time.time() - t0

    if result:
        flags = result.get('flags', [])
        contradictions = [f for f in flags if f.get('type') == 'contradiction']
        missing = [f for f in flags if f.get('type') == 'missing']
        confirmed = [f for f in flags if f.get('type') == 'confirmed']
        inferred = [f for f in flags if f.get('type') == 'inferred']

        print(f"  Validation complete ({elapsed:.1f}s):")
        print(f"    Confirmed: {len(confirmed)} fields")
        if inferred:
            print(f"    Inferred: {len(inferred)} new fields")
        if missing:
            print(f"    Missing context: {len(missing)} suggestions")
        if contradictions:
            print(f"    ⚠ Contradictions: {len(contradictions)} — needs review")
            for c in contradictions:
                print(f"      → {c.get('message', '')}")
    else:
        result = {"flags": [], "confidence": "low", "notes": "Validation failed — LLM returned invalid JSON"}

    result['_usage'] = {"input_tokens": usage.input_tokens, "output_tokens": usage.output_tokens}
    return result


# ══════════════════════════════════════════════════════════════
# MAIN: BUILD & SAVE CLIENT CONTEXT
# ══════════════════════════════════════════════════════════════

def build_context(
    tier1=None,
    tier2_answers=None,
    tier3_open=None,
    extract_from_files=None,
    extract_from_text=None,
    data_folder=None,
    output_path=None
):
    """
    Build a complete client_context.json from available inputs.

    Args:
        tier1: dict with Tier 1 fields (company_name, industry, etc.)
        tier2_answers: dict with answers to Tier 2 follow-up questions
        tier3_open: dict with open-ended context fields
        extract_from_files: list of file paths to extract context from (.pptx, .docx, etc.)
        extract_from_text: raw text (meeting notes, email, etc.) to extract from
        data_folder: path to client's data folder for cross-validation
        output_path: where to save client_context.json

    Returns:
        Complete context dict, validation result
    """
    print(f"\n{'='*70}")
    print("FINTASTIQ CLIENT CONTEXT INTAKE")
    print(f"{'='*70}")
    print(f"Run: {datetime.now().isoformat()}")

    context = {}
    provenance = {}
    total_in, total_out = 0, 0

    # ── Step 1: Start with Tier 1 fields ──
    if tier1:
        print(f"\nTier 1: Loading {len(tier1)} structured fields...")
        for k, v in tier1.items():
            if v is not None and v != '':
                context[k] = v
                provenance[k] = "team-entered"

    # ── Step 2: Extract from files if provided ──
    if extract_from_files:
        print(f"\nExtracting context from {len(extract_from_files)} file(s)...")
        for fpath in extract_from_files:
            extracted, usage = extract_context_from_file(fpath, context)
            if usage:
                total_in += usage.input_tokens
                total_out += usage.output_tokens
            if extracted:
                for k, v in extracted.items():
                    if v and k not in context:  # Don't overwrite explicit entries
                        context[k] = v
                        provenance[k] = f"extracted-from-{os.path.basename(fpath)}"
                    elif v and k in context:
                        # Store as supplementary
                        context.setdefault('_extracted_alternatives', {})[k] = {
                            "value": v,
                            "source": os.path.basename(fpath)
                        }

    # ── Step 3: Extract from raw text ──
    if extract_from_text:
        print(f"\nExtracting context from raw text ({len(extract_from_text):,} chars)...")
        extracted, usage = extract_context_from_text(extract_from_text, context)
        if usage:
            total_in += usage.input_tokens
            total_out += usage.output_tokens
        if extracted:
            for k, v in extracted.items():
                if v and k not in context:
                    context[k] = v
                    provenance[k] = "extracted-from-text"

    # ── Step 4: Validate Tier 1 completeness ──
    missing = validate_tier1(context)
    if missing:
        print(f"\n  ⚠ Missing required Tier 1 fields: {', '.join(missing)}")
    else:
        print(f"\n  ✓ All Tier 1 fields present")

    # ── Step 5: Generate Tier 2 follow-up questions ──
    if not missing:  # Only generate if Tier 1 is complete
        print(f"\nTier 2: Generating follow-up questions...")
        questions, usage = generate_tier2_questions(context)
        total_in += usage.input_tokens
        total_out += usage.output_tokens
        context['tier2_questions'] = questions

        # If answers were provided, attach them
        if tier2_answers:
            context['tier2_answers'] = tier2_answers
            for k, v in tier2_answers.items():
                provenance[f"tier2_{k}"] = "team-entered"

    # ── Step 6: Tier 3 open-ended context ──
    if tier3_open:
        print(f"\nTier 3: Loading {len(tier3_open)} open context fields...")
        context['open_context'] = tier3_open
        for k in tier3_open:
            provenance[f"open_{k}"] = "team-entered"

    # ── Step 7: Cross-validate against data ──
    validation_result = None
    if data_folder and os.path.isdir(data_folder):
        print(f"\nCross-validation: Checking context against data...")
        validation_result = cross_validate_context(context, data_folder)
        if validation_result:
            total_in += validation_result.get('_usage', {}).get('input_tokens', 0)
            total_out += validation_result.get('_usage', {}).get('output_tokens', 0)

            # Auto-apply inferred fields
            inferred = validation_result.get('inferred_fields', {})
            if inferred:
                for k, v in inferred.items():
                    if k not in context:
                        context[k] = v
                        provenance[k] = "inferred-from-data"
                print(f"  Auto-applied {len(inferred)} inferred fields")

            context['validation_flags'] = validation_result.get('flags', [])

    # ── Finalize ──
    context['provenance'] = provenance
    context['validated'] = validation_result.get('confidence', 'low') if validation_result else 'not-validated'
    context['created_at'] = datetime.now().isoformat()
    context['_meta'] = {
        "version": "1.0",
        "api_tokens": {"input": total_in, "output": total_out}
    }

    # Remove internal keys from output
    clean_context = {k: v for k, v in context.items() if not k.startswith('_') or k in ('_meta',)}

    # ── Save ──
    if output_path:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'w') as f:
            json.dump(clean_context, f, indent=2, default=str)
        print(f"\n  Saved: {output_path}")

    # ── Summary ──
    print(f"\n{'='*70}")
    print("CONTEXT INTAKE SUMMARY")
    print(f"{'='*70}")
    print(f"  Tier 1 fields: {sum(1 for k in TIER1_SCHEMA if k in context)}/{len(TIER1_SCHEMA)}")
    print(f"  Tier 2 questions: {len(context.get('tier2_questions', []))}")
    tier2_answered = len(context.get('tier2_answers', {}))
    if tier2_answered:
        print(f"  Tier 2 answered: {tier2_answered}")
    print(f"  Open context fields: {len(context.get('open_context', {}))}")
    print(f"  Provenance tracked: {len(provenance)} fields")
    if validation_result:
        flags = validation_result.get('flags', [])
        print(f"  Validation: {validation_result.get('confidence', '?')} confidence, {len(flags)} flags")
    print(f"  API tokens: {total_in:,} in / {total_out:,} out")

    if missing:
        print(f"\n  ⚠ INCOMPLETE — missing: {', '.join(missing)}")
        print(f"    Pipeline can run but with reduced LLM accuracy")
    else:
        print(f"\n  ✓ COMPLETE — ready to feed into pipeline")

    return clean_context, validation_result


# ══════════════════════════════════════════════════════════════
# HELPER: LOAD CONTEXT FOR PIPELINE USE
# ══════════════════════════════════════════════════════════════

def load_context(context_path):
    """Load client_context.json for use in pipeline stages."""
    if not os.path.exists(context_path):
        return None
    with open(context_path, 'r') as f:
        return json.load(f)


def context_to_prompt_block(context):
    """Convert context dict to a prompt-ready text block for LLM injection."""
    if not context:
        return ""

    lines = ["=== CLIENT CONTEXT ==="]

    # Core fields
    if context.get('company_name'):
        lines.append(f"Company: {context['company_name']}")
    if context.get('industry'):
        lines.append(f"Industry: {context['industry']}")
    if context.get('business_model'):
        lines.append(f"Business Model: {context['business_model']}")
    if context.get('revenue_range'):
        lines.append(f"Revenue Range: {context['revenue_range']}")
    if context.get('product_type'):
        pt = context['product_type']
        if isinstance(pt, list):
            pt = ', '.join(pt)
        lines.append(f"Products/Services: {pt}")
    if context.get('engagement_goal'):
        lines.append(f"Goal: {context['engagement_goal']}")

    # Key terms
    if context.get('key_terms') and isinstance(context['key_terms'], dict):
        lines.append("\nKey Terms:")
        for term, defn in context['key_terms'].items():
            lines.append(f"  {term} = {defn}")

    # Open context
    open_ctx = context.get('open_context', {})
    if open_ctx.get('pricing_structure'):
        lines.append(f"\nPricing Structure: {open_ctx['pricing_structure']}")
    if open_ctx.get('data_quirks'):
        lines.append(f"Data Quirks: {open_ctx['data_quirks']}")
    if open_ctx.get('known_issues'):
        lines.append(f"Known Issues: {open_ctx['known_issues']}")

    # Tier 2 answers
    tier2_answers = context.get('tier2_answers', {})
    if tier2_answers:
        lines.append("\nAdditional Context:")
        for qid, answer in tier2_answers.items():
            lines.append(f"  {qid}: {answer}")

    # Validation notes
    additional = context.get('additional_context')
    if additional:
        lines.append(f"\nAdditional: {additional}")

    lines.append("=== END CLIENT CONTEXT ===")
    return "\n".join(lines)


# ══════════════════════════════════════════════════════════════
# ENTRY POINT: NPI RETROFIT
# ══════════════════════════════════════════════════════════════

if __name__ == '__main__':
    # Retrofit NPI with context from what we know

    npi_tier1 = {
        "company_name": "NPI",
        "industry": "B2B SaaS / Financial Consulting",
        "business_model": "B2B",
        "revenue_range": "$10M-$50M",
        "product_type": ["Professional Services", "Subscriptions", "Software/SaaS"],
        "engagement_goal": "Pricing optimization, win rate improvement, sales funnel analysis"
    }

    npi_tier3 = {
        "pricing_structure": "Mix of project-based fees, subscription recurring revenue, audit services, and telecom. ACV ranges from <$10K to $100K+.",
        "known_issues": "Revenue gap between transaction fact ($119M) and monthly financials ($105M). Some data from 2022-2024 with gaps.",
        "data_quirks": "Source files include .xlsb binary workbooks. 'Bookings Dataset 1' has Excel serial dates. '0x2a' placeholder values in some fields. Account names have multiple variants (e.g., 'Aegon/Transamerica' vs 'Aegon Transamerica').",
        "key_terms": "TCV = Total Contract Value, ACV = Annual Contract Value, F500/F1000 = Fortune 500/1000 companies, ICP = Ideal Customer Profile"
    }

    data_folder = '/sessions/eloquent-laughing-hawking/mnt/Project Pricing Diagnostics/2. Client Internal Data/NPI Client Internal Data'
    output_path = '/sessions/eloquent-laughing-hawking/mnt/Project Pricing Diagnostics/2. Client Internal Data/NPI Client Internal Data/client_context.json'

    context, validation = build_context(
        tier1=npi_tier1,
        tier3_open=npi_tier3,
        data_folder=data_folder,
        output_path=output_path
    )
